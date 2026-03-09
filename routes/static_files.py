"""
routes/static_files.py — Static Asset Serving Blueprint (P2-T8)

Extracted from server.py during Phase 2 blueprint split.
Registers routes:
  GET  /sounds/<path:filepath>         — sound effect files
  GET  /uploads/<path:filename>        — uploaded user files
  GET  /src/<path:filepath>            — frontend JS/CSS source modules
  GET  /known_faces/<name>/<filename>  — face recognition photos
  GET  /api/dj-sound                   — DJ soundboard API (list/play)
"""

import logging
import random
import re
from pathlib import Path

from flask import Blueprint, Response, jsonify, request, send_file

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Blueprint
# ---------------------------------------------------------------------------

static_files_bp = Blueprint('static_files', __name__)

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

from services.paths import APP_ROOT, SOUNDS_DIR, UPLOADS_DIR, KNOWN_FACES_DIR, STATIC_DIR

BASE_DIR = APP_ROOT

UPLOADS_DIR.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------------------
# DJ Sounds catalogue
# ---------------------------------------------------------------------------

DJ_SOUNDS = {
    'air_horn': {
        'description': 'Classic stadium air horn - ba ba baaaa!',
        'when_to_use': 'Before drops, hype moments, celebrating wins, hip-hop DJ style'
    },
    'scratch_long': {
        'description': 'Extended DJ scratch solo - wicka wicka',
        'when_to_use': 'Transitions, hip-hop moments, showing off DJ skills'
    },
    'rewind': {
        'description': 'DJ rewind - pull up selecta!',
        'when_to_use': 'Going back, replaying something fire, dancehall pull-ups'
    },
    'record_stop': {
        'description': 'Record stopping abruptly',
        'when_to_use': 'Stopping everything, dramatic pause, cutting the music'
    },
    'impact': {
        'description': 'Punchy cinematic impact hit',
        'when_to_use': 'Punctuating statements, transitions, emphasis'
    },
    'crowd_cheer': {
        'description': 'Nightclub crowd cheering and going wild',
        'when_to_use': 'Big wins, amazing moments, festival energy, applause'
    },
    'crowd_hype': {
        'description': 'Hyped up rave crowd losing their minds',
        'when_to_use': 'Peak energy moments, party atmosphere'
    },
    'yeah': {
        'description': 'Hype man YEAH! vocal shot',
        'when_to_use': 'Hyping up, agreement, energy boost'
    },
    'lets_go': {
        'description': 'LETS GO! vocal chant',
        'when_to_use': 'Starting something, getting pumped, motivation'
    },
    'laser': {
        'description': 'Retro arcade laser zap - pew pew',
        'when_to_use': 'Sci-fi moments, gaming references, 80s vibes'
    },
    'gunshot': {
        'description': 'Dancehall gunshot sound - gun finger!',
        'when_to_use': 'Reggae/dancehall vibes, shooting down bad ideas'
    },
    'bruh': {
        'description': 'Classic bruh sound effect',
        'when_to_use': 'Facepalm moments, disappointment, when someone says something dumb'
    },
    'sad_trombone': {
        'description': 'Sad trombone wah wah wah - womp womp',
        'when_to_use': 'Fails, disappointments, when things go wrong'
    },
}

# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

def _safe_path(base_dir: Path, *parts) -> Path | None:
    """
    Resolve a path within base_dir, rejecting any traversal outside it.
    Returns the resolved Path on success, or None if traversal is detected.
    """
    try:
        resolved = (base_dir / Path(*parts)).resolve()
        base_resolved = base_dir.resolve()
        if resolved == base_resolved or base_resolved in resolved.parents:
            return resolved
    except Exception:
        pass
    return None


@static_files_bp.route('/sounds/<path:filepath>')
def serve_sound(filepath):
    """Serve sound effect files (including subdirectories like DJ-clips/)"""
    sound_path = _safe_path(SOUNDS_DIR, filepath)
    if sound_path is None:
        return jsonify({"error": "Invalid path"}), 400
    if sound_path.exists():
        return send_file(sound_path, mimetype='audio/mpeg')
    return jsonify({"error": "Sound not found"}), 404


# ---------------------------------------------------------------------------
# Upload constants & helpers
# ---------------------------------------------------------------------------

# Hard limit enforced before writing to disk (25 MB)
_MAX_UPLOAD_BYTES = 25 * 1024 * 1024

# Maximum characters returned to the AI as content_preview
_MAX_PREVIEW_CHARS = 6000

# Server-side allowlist — only these extensions are accepted
_ALLOWED_EXTENSIONS = {
    # Images
    '.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.tiff',
    # Structured documents
    '.pdf', '.docx', '.xlsx', '.pptx',
    # Plain text / code
    '.txt', '.md', '.csv', '.log',
    '.py', '.js', '.ts', '.json', '.yaml', '.yml',
    '.html', '.css',
}

# Control characters to strip from extracted text (keeps \t \n \r)
_CTRL_RE = re.compile(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]')


def _sanitize_text(text: str) -> str:
    """Strip control chars, collapse excessive blank lines, cap at _MAX_PREVIEW_CHARS."""
    text = _CTRL_RE.sub('', text)
    text = re.sub(r'\n{4,}', '\n\n\n', text)  # no more than 3 consecutive blank lines
    return text[:_MAX_PREVIEW_CHARS].strip()


def _extract_pdf(path: Path) -> str:
    """Extract text from a PDF using pypdf. Returns sanitized string."""
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ''
            pages.append(text)
        except Exception:
            pages.append(f'[Page {i + 1}: extraction failed]')
    return _sanitize_text('\n\n'.join(pages))


def _extract_docx(path: Path) -> str:
    """Extract text from a .docx using python-docx. Returns sanitized string."""
    from docx import Document
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(' | '.join(cells))
    return _sanitize_text('\n'.join(parts))


def _extract_xlsx(path: Path) -> str:
    """Extract cell values from a .xlsx using openpyxl. Returns sanitized string."""
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    try:
        for sheet in wb.worksheets:
            parts.append(f'[Sheet: {sheet.title}]')
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    parts.append('\t'.join(cells))
    finally:
        wb.close()
    return _sanitize_text('\n'.join(parts))


def _extract_pptx(path: Path) -> str:
    """Extract text from a .pptx using python-pptx. Returns sanitized string."""
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            parts.append(f'[Slide {i}]')
            parts.extend(slide_texts)
    return _sanitize_text('\n'.join(parts))


def _call_skill_runner(path: Path, original_name: str) -> str | None:
    """
    Try to extract document text via the shared skill-runner service.
    Returns extracted text on success, None if the service is unavailable.
    Falls back gracefully so local extractors can take over.
    """
    try:
        import requests
        with open(path, 'rb') as fh:
            resp = requests.post(
                'http://skill-runner:8900/extract',
                files={'file': (original_name, fh)},
                data={'filename': original_name},
                timeout=30,
            )
        if resp.ok:
            data = resp.json()
            return data.get('text', '')
        logger.warning('skill-runner /extract returned %d for %s', resp.status_code, original_name)
    except Exception as exc:
        logger.debug('skill-runner unavailable, using local extractors: %s', exc)
    return None


@static_files_bp.route('/api/upload', methods=['POST'])
def upload_file():
    """Accept a file upload from the text panel and save to uploads/."""
    import mimetypes
    import uuid

    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400

    f = request.files['file']
    if not f.filename:
        return jsonify({'error': 'Empty filename'}), 400

    # --- Sanitize filename, validate extension ---
    original_name = Path(f.filename).name
    ext = Path(original_name).suffix.lower()
    if ext not in _ALLOWED_EXTENSIONS:
        return jsonify({'error': f'File type "{ext}" is not allowed'}), 415

    # --- Size check before writing to disk ---
    # Seek to end to get byte length without reading into memory
    f.stream.seek(0, 2)
    file_size = f.stream.tell()
    f.stream.seek(0)
    if file_size > _MAX_UPLOAD_BYTES:
        return jsonify({'error': 'File too large (25 MB max)'}), 413

    # --- Save with UUID filename (no original name on disk) ---
    safe_name = f"{uuid.uuid4().hex}{ext}"
    dest = UPLOADS_DIR / safe_name
    UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
    f.save(str(dest))

    mime = f.mimetype or mimetypes.guess_type(original_name)[0] or ''
    is_image = mime.startswith('image/')

    result = {
        'original_name': original_name,
        'path': str(dest),
        'filename': safe_name,
        'url': f'/uploads/{safe_name}',
    }

    if is_image:
        result['type'] = 'image'
        return jsonify(result)

    result['type'] = 'file'

    # --- Extract readable content by type ---
    _BINARY_EXTS = {'.pdf', '.docx', '.xlsx', '.pptx'}

    try:
        if ext in _BINARY_EXTS:
            # Try shared skill-runner first (preferred — keeps this container lean)
            text = _call_skill_runner(dest, original_name)

            # Fall back to local extractors if skill-runner unavailable
            if text is None:
                if ext == '.pdf':
                    text = _extract_pdf(dest)
                elif ext == '.docx':
                    text = _extract_docx(dest)
                elif ext == '.xlsx':
                    text = _extract_xlsx(dest)
                elif ext == '.pptx':
                    text = _extract_pptx(dest)

            if text:
                result['content_preview'] = text
                result['extracted_type'] = ext.lstrip('.')
            else:
                result['extraction_error'] = (
                    f'Could not extract text from {ext} file. '
                    'Install skill-runner or document packages to enable this.'
                )

        else:
            # Plain text / code / CSV — read directly
            text_types = {'text/', 'application/json', 'application/xml', 'application/javascript'}
            if any(mime.startswith(t) for t in text_types) or ext in {
                '.txt', '.md', '.csv', '.log',
                '.py', '.js', '.ts', '.json', '.yaml', '.yml',
                '.html', '.css',
            }:
                raw = dest.read_text(errors='replace')
                result['content_preview'] = _sanitize_text(raw)

    except Exception as exc:
        logger.warning('Document extraction failed for %s: %s', original_name, exc)
        result['extraction_error'] = f'Could not extract text from {ext} file'

    return jsonify(result)


@static_files_bp.route('/static/emulators/<path:filepath>')
def serve_emulator(filepath):
    """Serve bundled emulator files (js-dos, etc.) from /app/static/emulators/."""
    emulators_dir = STATIC_DIR / 'emulators'
    path = _safe_path(emulators_dir, filepath)
    if path is None:
        return jsonify({"error": "Invalid path"}), 400
    if not path.exists():
        return jsonify({"error": "File not found"}), 404
    mime_types = {'.js': 'application/javascript', '.css': 'text/css', '.wasm': 'application/wasm'}
    mime = mime_types.get(path.suffix, 'application/octet-stream')
    response = send_file(path, mimetype=mime)
    response.headers['Cache-Control'] = 'public, max-age=86400'
    response.headers['Access-Control-Allow-Origin'] = '*'
    return response


@static_files_bp.route('/uploads/<path:filename>')
def serve_upload(filename):
    """Serve uploaded files (path traversal guarded)."""
    upload_path = _safe_path(UPLOADS_DIR, filename)
    if upload_path is None:
        return jsonify({"error": "Invalid path"}), 400
    if not upload_path.exists():
        return jsonify({"error": "File not found"}), 404
    return send_file(upload_path)


@static_files_bp.route('/src/<path:filepath>')
def serve_src(filepath):
    """Serve frontend source files (JS, CSS modules)"""
    # P7-T3 security: prevent path traversal (same guard used by serve_sound)
    src_path = _safe_path(APP_ROOT / 'src', filepath)
    if src_path is None:
        return jsonify({"error": "Invalid path"}), 400
    if not src_path.exists():
        return jsonify({"error": "File not found"}), 404

    mime_types = {
        '.js': 'application/javascript',
        '.css': 'text/css',
        '.html': 'text/html',
        '.json': 'application/json',
    }
    mime_type = mime_types.get(src_path.suffix.lower(), 'text/plain')
    resp = send_file(src_path, mimetype=mime_type)
    resp.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
    return resp


@static_files_bp.route('/known_faces/<name>/<filename>')
def serve_face_photo(name, filename):
    """Serve face photos for the My Face section"""
    photo_path = _safe_path(KNOWN_FACES_DIR, name, filename)
    if photo_path is None:
        return jsonify({"error": "Invalid path"}), 400
    if photo_path.exists():
        return send_file(photo_path)
    return jsonify({"error": "Photo not found"}), 404


@static_files_bp.route('/api/dj-sound', methods=['GET'])
def handle_dj_sound():
    """
    DJ Soundboard endpoint.
    Query params:
      - action: 'list' or 'play'
      - sound: sound name (e.g., 'air_horn', 'scratch', 'siren_rise')
    Returns sound info or triggers playback.
    """
    action = request.args.get('action', 'list')
    sound = request.args.get('sound', '')

    if action == 'list':
        sounds_list = [
            {
                'name': name,
                'description': info['description'],
                'when_to_use': info['when_to_use'],
                'available': (SOUNDS_DIR / f"{name}.mp3").exists(),
            }
            for name, info in DJ_SOUNDS.items()
        ]
        return jsonify({
            'action': 'list',
            'sounds': sounds_list,
            'count': len(sounds_list),
            'response': (
                f"Soundboard loaded! {len(sounds_list)} effects ready. "
                "I got air horns, sirens, scratches, crowd effects, and more!"
            ),
        })

    if action == 'play':
        if not sound:
            sound = random.choice(list(DJ_SOUNDS.keys()))

        sound_lower = sound.lower().replace(' ', '_').replace('-', '_')

        matched = next(
            (name for name in DJ_SOUNDS if sound_lower in name or name in sound_lower),
            None,
        )
        if not matched:
            matched = next(
                (name for name in DJ_SOUNDS
                 if any(word in name for word in sound_lower.split('_'))),
                None,
            )

        if not matched:
            return jsonify({
                'action': 'error',
                'response': (
                    f"No sound matching '{sound}'. "
                    "Try: air_horn, siren, scratch, applause, bass_drop, rewind..."
                ),
            })

        sound_file = SOUNDS_DIR / f"{matched}.mp3"
        if not sound_file.exists():
            return jsonify({
                'action': 'error',
                'response': f"Sound file for '{matched}' not found. Need to generate it first!",
            })

        info = DJ_SOUNDS[matched]
        return jsonify({
            'action': 'play',
            'sound': matched,
            'description': info['description'],
            'url': f"/sounds/{matched}.mp3",
            'response': f"*{info['description'].upper()}* 🎵",
        })

    return jsonify({'error': 'Unknown action'}), 400


@static_files_bp.route('/manifest.json')
def serve_manifest():
    """PWA Web App Manifest — dynamically injects CLIENT_NAME for per-tenant PWA identity"""
    import json as _json, os as _os
    client_name = _os.environ.get("CLIENT_NAME", "").strip()
    path = STATIC_DIR / 'manifest.json'
    manifest = _json.loads(path.read_text())
    if client_name:
        manifest["name"] = client_name
        manifest["short_name"] = client_name
    resp = Response(
        _json.dumps(manifest, indent=2),
        mimetype='application/manifest+json'
    )
    resp.headers['Cache-Control'] = 'public, max-age=86400'
    return resp


@static_files_bp.route('/sw.js')
def serve_sw():
    """PWA Service Worker — must be served from root scope"""
    path = STATIC_DIR / 'sw.js'
    resp = send_file(path, mimetype='application/javascript')
    resp.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
    resp.headers['Service-Worker-Allowed'] = '/'
    return resp


@static_files_bp.route('/static/icons/<filename>')
def serve_icon(filename):
    """PWA icons"""
    icon_path = _safe_path(STATIC_DIR / 'icons', filename)
    if icon_path is None or not icon_path.exists():
        return jsonify({"error": "Icon not found"}), 404
    return send_file(icon_path, mimetype='image/png')


@static_files_bp.route('/install')
def serve_install():
    """PWA install landing page"""
    path = STATIC_DIR / 'install.html'
    resp = send_file(path, mimetype='text/html')
    resp.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
    return resp


@static_files_bp.route('/admin')
def serve_admin():
    """Serve the OpenUI admin dashboard"""
    admin_path = APP_ROOT / 'src' / 'admin.html'
    if not admin_path.exists():
        return jsonify({"error": "Admin dashboard not found"}), 404
    resp = send_file(admin_path, mimetype='text/html')
    resp.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
    return resp
